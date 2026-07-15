"""PPTX (PowerPoint) fayllarni ingest qilish (python-pptx)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ufl.ingest.base import Block, Document


def extract(path: Path) -> Document:
    presentation = Presentation(str(path))
    blocks = []
    for slide_index, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                blocks.append(Block(text=text, page=slide_index))
    return Document(blocks=blocks)
