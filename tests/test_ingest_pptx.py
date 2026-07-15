from pptx import Presentation
from pptx.util import Inches

from ufl.ingest.pptx import extract


def test_pptx_extract_reads_nonempty_slide_text(tmp_path):
    presentation = Presentation()
    layout = presentation.slide_layouts[6]  # bo'sh layout

    slide1 = presentation.slides.add_slide(layout)
    box1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box1.text_frame.text = "Birinchi slayd matni."

    slide2 = presentation.slides.add_slide(layout)
    box2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box2.text_frame.text = ""  # bo'sh matn o'tkazib yuborilishi kerak
    box3 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box3.text_frame.text = "Ikkinchi slayd matni."

    path = tmp_path / "sample.pptx"
    presentation.save(str(path))

    document = extract(path)

    assert [b.text for b in document.blocks] == ["Birinchi slayd matni.", "Ikkinchi slayd matni."]
    assert document.blocks[0].page == 0
    assert document.blocks[1].page == 1


def test_detect_format_recognizes_pptx(tmp_path):
    from ufl.ingest.detect import detect_format

    presentation = Presentation()
    path = tmp_path / "sample.pptx"
    presentation.save(str(path))

    assert detect_format(path) == "pptx"
